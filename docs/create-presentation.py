#!/usr/bin/env python3
"""Generate PowerPoint presentation for Assistant Framework v0.2.0"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SOFT_BG = RGBColor(0xF5, 0xF5, 0xF5)
CARD_BG = RGBColor(0x24, 0x24, 0x3E)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT = RGBColor(0xFF, 0x45, 0x45)
GREEN = RGBColor(0x00, 0xE6, 0x76)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_dark_bg(slide):
    """Add dark background to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_light_bg(slide):
    """Add light background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SOFT_BG


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide)

add_text_box(slide, 1, 1.5, 11, 1.2, "ASSISTANT FRAMEWORK", 48, ACCENT, True)
add_accent_line(slide, 1, 2.7, 4)
add_text_box(slide, 1, 3.0, 11, 1.5, "Your AI Becomes a Senior Developer\nThat Never Forgets", 32, WHITE, False)
add_text_box(slide, 1, 5.0, 11, 0.5, "v0.2.0  |  11 Skills  |  13 MCP Tools  |  Self-Improving", 18, MED_GRAY)
add_text_box(slide, 1, 6.0, 11, 0.5, "Claude Code  ·  OpenAI Codex  ·  Google Gemini CLI", 16, MED_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "THE PROBLEM", 14, ACCENT, True)
add_text_box(slide, 1, 1.2, 11, 1, "Every AI coding session starts from zero.", 36, WHITE, True)
add_accent_line(slide, 1, 2.3, 3)

problems = [
    "No memory of past decisions or what worked before",
    "Repeats the same mistakes across sessions",
    "Doesn't know your preferences or coding style",
    "No structured workflow — just vibes and hope",
    "Can't cover your weaknesses (docs, diagrams, onboarding)",
    "No quality gates — ships whatever comes out first",
]
add_bullet_list(slide, 1.2, 2.8, 10, 3.5, problems, 20, LIGHT_GRAY, Pt(14))

add_card(slide, 1, 6.0, 11.3, 0.8, RGBColor(0x3A, 0x1A, 0x1A))
add_text_box(slide, 1.3, 6.1, 10.7, 0.6,
             "You are the memory. You are the process. You are the quality gate.", 18, RED_ACCENT, True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: The Solution Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "THE SOLUTION", 14, ACCENT, True)
add_text_box(slide, 1, 1.2, 11, 1, "One framework. Self-improving. Cross-platform.", 32, WHITE, True)
add_accent_line(slide, 1, 2.2, 3)

# Skill cards - row 1
cards_r1 = [
    ("Workflow", "Structured dev\npipeline with\napproval gates"),
    ("Review", "Autonomous loop\nmax 5 rounds\nfresh reviewer"),
    ("Security", "STRIDE + OWASP\nCVE audit\nattack surface"),
    ("TDD", "Red-Green-Refactor\nstrict gates\nno code before test"),
]

for i, (title, desc) in enumerate(cards_r1):
    x = 0.8 + i * 3.1
    add_card(slide, x, 2.8, 2.8, 1.8)
    add_text_box(slide, x + 0.2, 2.9, 2.4, 0.4, title, 16, ACCENT, True)
    add_text_box(slide, x + 0.2, 3.3, 2.4, 1.2, desc, 13, LIGHT_GRAY)

# Skill cards - row 2
cards_r2 = [
    ("Docs", "6 modes: API, arch\nREADME, changelog\nmigration, explain"),
    ("Onboard", "6-phase codebase\nlearning protocol\nauto-generates memory"),
    ("Ideate", "Diverge-converge\nrefine pipeline\nscored ranking"),
    ("Reflexion", "Self-improving\nlesson recall\nstrategy profiles"),
]

for i, (title, desc) in enumerate(cards_r2):
    x = 0.8 + i * 3.1
    add_card(slide, x, 4.9, 2.8, 1.8)
    add_text_box(slide, x + 0.2, 5.0, 2.4, 0.4, title, 16, ACCENT, True)
    add_text_box(slide, x + 0.2, 5.4, 2.4, 1.2, desc, 13, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Structured Workflow
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "STRUCTURED WORKFLOW", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Right-sized ceremony. Invisible when working.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

# Pipeline visualization
phases = ["TRIAGE", "DISCOVER", "PLAN", "BUILD", "REVIEW", "DOCUMENT"]
for i, phase in enumerate(phases):
    x = 0.5 + i * 2.1
    add_card(slide, x, 2.5, 1.9, 0.7, ACCENT2)
    add_text_box(slide, x, 2.55, 1.9, 0.6, phase, 13, WHITE, True, PP_ALIGN.CENTER)
    if i < len(phases) - 1:
        add_text_box(slide, x + 1.9, 2.55, 0.3, 0.6, "→", 18, ACCENT, True, PP_ALIGN.CENTER)

# Size table
add_text_box(slide, 1, 3.6, 5, 0.4, "Adapts to task size:", 16, ACCENT, True)

sizes = [
    ("Small", "bugfix, typo", "Writer → Tester → Reviewer"),
    ("Medium", "feature, refactor", "Mapper → Writer → Tester → Reviewer"),
    ("Large", "new project", "Mapper → Explorer → Architect → Writer → Tester → Reviewer"),
    ("Mega", "rewrite, 10+ files", "Full pipeline + parallel sub-tasks"),
]
for i, (size, desc, flow) in enumerate(sizes):
    y = 4.1 + i * 0.55
    add_text_box(slide, 1.2, y, 1.5, 0.5, size, 15, ORANGE, True)
    add_text_box(slide, 2.7, y, 2.3, 0.5, desc, 14, MED_GRAY)
    add_text_box(slide, 5.0, y, 7.5, 0.5, flow, 14, LIGHT_GRAY)

add_card(slide, 1, 6.3, 11.3, 0.7, RGBColor(0x1A, 0x2E, 0x1A))
add_text_box(slide, 1.3, 6.35, 10.7, 0.6,
             "Key principle: if the user notices the framework, it's too heavy.", 16, GREEN, True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Autonomous Code Review
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "AUTONOMOUS CODE REVIEW", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Not one pass. An autonomous fix-and-review loop.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

rounds = [
    ("Round 1", "80%+ confidence", "Fresh Reviewer finds 4 issues → Fix all → Tests pass", ORANGE),
    ("Round 2", "85%+ confidence", "NEW Reviewer finds 1 more → Fix → Tests pass", ORANGE),
    ("Round 3", "90%+ confidence", "NEW Reviewer → Clean. No findings above nit.", GREEN),
]

for i, (rnd, conf, desc, clr) in enumerate(rounds):
    y = 2.4 + i * 1.2
    add_card(slide, 1, y, 11.3, 1.0)
    add_text_box(slide, 1.3, y + 0.05, 1.5, 0.4, rnd, 18, clr, True)
    add_text_box(slide, 3.0, y + 0.05, 2, 0.4, conf, 14, MED_GRAY)
    add_text_box(slide, 1.3, y + 0.5, 10.5, 0.4, desc, 16, LIGHT_GRAY)

features = [
    "Fresh reviewer each round — stale context weakens reviews",
    "Previously-fixed list prevents re-reporting same issues",
    "Stop hook structurally prevents finishing without review",
]
add_bullet_list(slide, 1.2, 6.2, 11, 1.2, features, 14, MED_GRAY, Pt(4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Documentation & Diagrams
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "DOCUMENTATION & DIAGRAMS", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Cover your biggest weakness automatically.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

# Left column - Docs
add_text_box(slide, 1, 2.4, 5, 0.5, "Documentation Generator", 20, ACCENT, True)
doc_modes = [
    "API Docs — scans endpoints, generates reference",
    "Architecture — system overview with diagrams",
    "README — from code analysis, not memory",
    "Changelog — from git history, categorized",
    "Migration Guide — breaking changes + steps",
    "Code Explainer — the 'why', not the 'what'",
]
add_bullet_list(slide, 1.2, 3.0, 5.5, 3, doc_modes, 15, LIGHT_GRAY, Pt(10))

# Right column - Diagrams
add_text_box(slide, 7, 2.4, 5, 0.5, "Diagram Generator", 20, ACCENT, True)
diagram_types = [
    "Architecture — component relationships",
    "Sequence — request flows, interactions",
    "Entity-Relationship — data models",
    "Flow — decision trees, algorithms",
    "Component — module dependencies",
    "Class — type hierarchies",
    "State — lifecycle transitions",
]
add_bullet_list(slide, 7.2, 3.0, 5.5, 3, diagram_types, 15, LIGHT_GRAY, Pt(10))

add_card(slide, 1, 5.8, 11.3, 1.2, RGBColor(0x1A, 0x1A, 0x3E))
add_text_box(slide, 1.3, 5.9, 10.7, 0.4,
             "Also detects stale docs:", 15, ACCENT, True)
add_text_box(slide, 1.3, 6.3, 10.7, 0.6,
             "\"README.md: last updated 45 days ago, 3 new features since\"", 16, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: The Reflexion Breakthrough
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "THE BREAKTHROUGH: REFLEXION", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Every task makes the next task better.", 32, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

# Timeline
sessions = [
    ("Session 1", "Fix API bug",
     "Reflexion: \"Wasted 5 min — forgot to check DI registration\"\n"
     "Lesson stored with confidence 0.5"),
    ("Session 5", "Another API bug",
     "Discover phase: \"Found 2 relevant lessons\"\n"
     "Applies lesson automatically → faster fix\n"
     "Lesson reinforced → confidence 0.7"),
    ("Session 20", "API refactor",
     "12 accumulated lessons across all phases\n"
     "\"You underestimate refactors by 1 size\" → auto-adjusts\n"
     "Strategy profile guides the entire approach"),
]

for i, (session, task, detail) in enumerate(sessions):
    y = 2.4 + i * 1.35
    clr = [ORANGE, ACCENT, GREEN][i]
    add_card(slide, 1, y, 11.3, 1.2)
    add_text_box(slide, 1.3, y + 0.05, 2, 0.4, session, 17, clr, True)
    add_text_box(slide, 3.5, y + 0.05, 3, 0.4, task, 15, WHITE)
    add_text_box(slide, 1.3, y + 0.4, 10.5, 0.7, detail, 12, LIGHT_GRAY)

add_text_box(slide, 1, 6.7, 11.3, 0.5,
             "Confidence scoring · Time decay · Automatic consolidation · Per-project-type strategies",
             14, MED_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Memory Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "MEMORY ARCHITECTURE", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Not just files. A queryable knowledge system.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

# Left side - Knowledge Graph
add_card(slide, 0.8, 2.4, 5.5, 4.5)
add_text_box(slide, 1.1, 2.5, 5, 0.4, "Knowledge Graph (JSONL)", 17, ACCENT, True)
graph_items = [
    "Entities: Projects, Technologies, Patterns",
    "Relations: DependsOn, Uses, Follows",
    "Insights linked to projects",
    "Conventions per codebase",
    "Preferences (global + scoped)",
    "Markdown sync on startup",
]
add_bullet_list(slide, 1.3, 3.0, 4.8, 3.5, graph_items, 14, LIGHT_GRAY, Pt(10))

# Right side - SQLite + FTS5
add_card(slide, 7, 2.4, 5.5, 4.5)
add_text_box(slide, 7.3, 2.5, 5, 0.4, "SQLite + FTS5 (v2)", 17, ACCENT, True)
sqlite_items = [
    "Reflexions: post-task self-assessments",
    "Decisions: with rationale + alternatives",
    "Strategy Lessons: per project type",
    "Calibration: prediction accuracy tracking",
    "FTS5 Index: ranked search across ALL content",
    "Porter stemming + Unicode tokenizer",
]
add_bullet_list(slide, 7.3, 3.0, 4.8, 3.5, sqlite_items, 14, LIGHT_GRAY, Pt(10))

# MCP tools bar
add_text_box(slide, 1, 7.0, 11.3, 0.4,
             "13 MCP tools: context · search · reflect · decide · pattern · consolidate · stats · add entity/relation/insight · remove · graph",
             13, MED_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Multi-Agent Orchestration
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "MULTI-AGENT ORCHESTRATION", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Specialized roles with constrained access.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

agents = [
    ("Code Mapper", "Read-only", "Lightweight structural map of the codebase"),
    ("Explorer", "Read-only", "Deep execution path tracing, hidden dependencies"),
    ("Architect", "Read-only", "Implementation blueprints, component design"),
    ("Code Writer", "Write", "Implements code following the plan — nothing else"),
    ("Builder/Tester", "Write", "Builds, writes tests, runs tests, absorbs noise"),
    ("Reviewer", "Read-only", "Independent review with confidence filtering"),
]

for i, (name, access, desc) in enumerate(agents):
    y = 2.4 + i * 0.75
    access_color = GREEN if access == "Read-only" else ORANGE
    add_text_box(slide, 1.2, y, 2.5, 0.5, name, 16, WHITE, True)
    add_text_box(slide, 3.8, y, 1.8, 0.5, access, 14, access_color, True)
    add_text_box(slide, 5.8, y, 6.5, 0.5, desc, 14, LIGHT_GRAY)

add_card(slide, 1, 6.2, 11.3, 0.8, RGBColor(0x1A, 0x1A, 0x3E))
add_text_box(slide, 1.3, 6.3, 10.7, 0.6,
             "Reviewer cannot edit files. Code Writer doesn't run tests.\n"
             "Separation of concerns at the agent level.",
             15, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Automated Hooks
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "AUTOMATED HOOKS", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Six lifecycle hooks. Zero manual steps.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

hooks = [
    ("Session Start", "New session begins", "Injects memory, task state, reflexion tools"),
    ("Skill Router", "Every user prompt", "Auto-routes to the correct skill"),
    ("Pre-Compress", "Before compaction", "Saves state before context is lost"),
    ("Post-Compact", "After compaction", "Re-injects task journal and rules"),
    ("Stop Review", "Agent tries to finish", "BLOCKS until review cycle is complete"),
    ("Session End", "Session closes", "Prompts for reflexion + memory capture"),
]

for i, (name, when, what) in enumerate(hooks):
    y = 2.4 + i * 0.7
    name_color = RED_ACCENT if name == "Stop Review" else ACCENT
    add_text_box(slide, 1.2, y, 2.5, 0.5, name, 15, name_color, True)
    add_text_box(slide, 3.8, y, 3, 0.5, when, 14, MED_GRAY)
    add_text_box(slide, 7, y, 5.5, 0.5, what, 14, LIGHT_GRAY)

add_card(slide, 1, 6.3, 11.3, 0.7, RGBColor(0x3A, 0x1A, 0x1A))
add_text_box(slide, 1.3, 6.35, 10.7, 0.6,
             "The stop-review hook is structural enforcement — the agent physically cannot finish without review.",
             15, RED_ACCENT, True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: What Makes This Different
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "WHAT MAKES THIS DIFFERENT", 14, ACCENT, True)
add_text_box(slide, 1, 1.1, 11, 0.7, "Compared to other frameworks.", 28, WHITE, True)
add_accent_line(slide, 1, 1.9, 3)

comparisons = [
    ("Others: Skills as suggestions", "Ours: Skills as mandatory enforcement"),
    ("Others: Memory as chat history", "Ours: Queryable knowledge graph + FTS5"),
    ("Others: One-shot review", "Ours: Autonomous loop, max 5 rounds"),
    ("Others: Generic patterns", "Ours: YOUR project's patterns, learned over time"),
    ("Others: Reactive only", "Ours: Self-improving via reflexion"),
    ("Others: Single platform", "Ours: Claude + Codex + Gemini"),
    ("Others: Import from marketplace", "Ours: 100% built in-house"),
]

for i, (them, us) in enumerate(comparisons):
    y = 2.4 + i * 0.68
    add_text_box(slide, 1.2, y, 5.5, 0.5, them, 14, MED_GRAY)
    add_text_box(slide, 7, y, 5.5, 0.5, us, 14, GREEN, True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: By The Numbers
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 0.6, 11, 0.8, "BY THE NUMBERS", 14, ACCENT, True)
add_accent_line(slide, 1, 1.3, 3)

numbers = [
    ("11", "Skills"),
    ("13", "MCP Tools"),
    ("6", "Lifecycle Hooks"),
    ("6", "Specialized Agents"),
    ("7", "Diagram Types"),
    ("6", "Doc Modes"),
]

for i, (num, label) in enumerate(numbers):
    col = i % 3
    row = i // 3
    x = 1.5 + col * 3.8
    y = 1.8 + row * 2.5
    add_text_box(slide, x, y, 3, 1, num, 64, ACCENT, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + 1.1, 3, 0.5, label, 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Bottom stats
add_text_box(slide, 1, 6.5, 11.3, 0.5,
             "91 tests passing  ·  1 external dependency  ·  0 marketplace imports  ·  3 platforms",
             16, MED_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Closing
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, 1, 2.0, 11.3, 1, "ASSISTANT FRAMEWORK", 44, ACCENT, True, PP_ALIGN.CENTER)
add_accent_line(slide, 5.5, 3.2, 2.3)
add_text_box(slide, 1, 3.5, 11.3, 1.5,
             "Your AI.  Your workflow.  Your memory.\n\nIt learns.  It improves.  It never forgets.",
             28, WHITE, False, PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.8, 11.3, 1,
             "git clone <repo> && ./install.sh --agent claude\n\nThat's it. Skills auto-trigger. Hooks auto-fire. Memory accumulates.",
             16, MED_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Assistant-Framework-v0.2.0.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
